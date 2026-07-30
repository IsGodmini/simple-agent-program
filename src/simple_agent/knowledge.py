"""Local document ingestion and retrieval for project knowledge RAG."""

import csv
import hashlib
import html
import io
import json
import re
import sqlite3
import zipfile
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Set, Tuple
from xml.etree import ElementTree

from .workspace import Workspace
from .storage import ProjectStorage
from .vector_store import (
    ChromaVectorStore,
    VectorRecord,
    content_fingerprint,
    reciprocal_rank_fusion,
)

KNOWLEDGE_VERSION = 1
DEFAULT_CHUNK_CHARS = 1_800
DEFAULT_CHUNK_OVERLAP = 200
MAX_SOURCE_BYTES = 50 * 1024 * 1024
MAX_EXTRACTED_CHARS = 2_000_000
MAX_ARCHIVE_BYTES = 200 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 10_000
VALID_DOCUMENT_ID = re.compile(r"^[a-f0-9]{16}$")

TEXT_EXTENSIONS = {
    ".cfg",
    ".conf",
    ".css",
    ".env.example",
    ".go",
    ".graphql",
    ".ini",
    ".java",
    ".js",
    ".jsx",
    ".kt",
    ".log",
    ".md",
    ".mdx",
    ".properties",
    ".py",
    ".rb",
    ".rs",
    ".rst",
    ".sh",
    ".sql",
    ".swift",
    ".toml",
    ".ts",
    ".tsx",
    ".txt",
    ".vue",
    ".yaml",
    ".yml",
}
STRUCTURED_EXTENSIONS = {
    ".csv",
    ".docx",
    ".epub",
    ".htm",
    ".html",
    ".json",
    ".jsonl",
    ".odt",
    ".pdf",
    ".ppsx",
    ".pptx",
    ".tsv",
    ".xlsm",
    ".xlsx",
    ".xml",
}


@dataclass(frozen=True)
class ParsedSection:
    """One logical section extracted from an uploaded document."""

    heading: str
    location: str
    text: str


@dataclass(frozen=True)
class KnowledgeDocument:
    """Metadata returned after one source document is indexed."""

    document_id: str
    source_name: str
    source_type: str
    content_hash: str
    chunk_count: int
    imported_at: str


@dataclass(frozen=True)
class KnowledgeHit:
    """One retrieved knowledge chunk with a stable citation."""

    document_id: str
    source_name: str
    chunk_index: int
    heading: str
    location: str
    content: str
    score: float

    @property
    def citation(self) -> str:
        return f"knowledge:{self.document_id}#chunk-{self.chunk_index}"


class _VisibleHTMLParser(HTMLParser):
    """Extract visible text without retaining script or style contents."""

    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []
        self._ignored_depth = 0

    def handle_starttag(
        self,
        tag: str,
        attrs: List[Tuple[str, Optional[str]]],
    ) -> None:
        if tag.lower() in {"script", "style"}:
            self._ignored_depth += 1
        elif not self._ignored_depth and tag.lower() in {
            "br",
            "div",
            "h1",
            "h2",
            "h3",
            "h4",
            "li",
            "p",
            "section",
            "tr",
        }:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style"} and self._ignored_depth:
            self._ignored_depth -= 1
        elif not self._ignored_depth and tag.lower() in {
            "div",
            "h1",
            "h2",
            "h3",
            "h4",
            "li",
            "p",
            "section",
            "tr",
        }:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth:
            self.parts.append(data)

    def text(self) -> str:
        return _normalize_text(html.unescape("".join(self.parts)))


class DocumentParser:
    """Parse common developer and office document formats into text sections."""

    def parse(self, path: Path) -> List[ParsedSection]:
        if not path.exists():
            raise ValueError(f"knowledge file does not exist: {path}")
        if not path.is_file():
            raise ValueError(f"knowledge path is not a file: {path}")
        size = path.stat().st_size
        if size > MAX_SOURCE_BYTES:
            raise ValueError(
                f"knowledge file exceeds {MAX_SOURCE_BYTES} bytes: {path.name}"
            )

        suffix = path.suffix.lower()
        if suffix in {
            ".docx",
            ".epub",
            ".odt",
            ".ppsx",
            ".pptx",
            ".xlsm",
            ".xlsx",
        }:
            _validate_archive(path)
        if suffix == ".pdf":
            sections = self._parse_pdf(path)
        elif suffix == ".docx":
            sections = self._parse_docx(path)
        elif suffix in {".pptx", ".ppsx"}:
            sections = self._parse_pptx(path)
        elif suffix in {".xlsx", ".xlsm"}:
            sections = self._parse_xlsx(path)
        elif suffix in {".csv", ".tsv"}:
            sections = self._parse_delimited(path, suffix)
        elif suffix == ".json":
            sections = self._parse_json(path)
        elif suffix == ".jsonl":
            sections = self._parse_jsonl(path)
        elif suffix in {".html", ".htm"}:
            sections = [
                ParsedSection(path.stem, "document", self._html_text(path))
            ]
        elif suffix == ".xml":
            sections = self._parse_xml(path)
        elif suffix == ".odt":
            sections = self._parse_odt(path)
        elif suffix == ".epub":
            sections = self._parse_epub(path)
        elif suffix in {".doc", ".xls", ".ppt"}:
            raise ValueError(
                f"legacy Office format {suffix} is not supported; "
                "please convert it to the corresponding x format"
            )
        else:
            sections = self._parse_text(path)

        cleaned = self._bounded_sections(sections)
        if not cleaned:
            raise ValueError(f"no readable text was extracted from: {path.name}")
        return cleaned

    def _parse_text(self, path: Path) -> List[ParsedSection]:
        raw = path.read_bytes()
        if b"\x00" in raw[:8192]:
            raise ValueError(f"unsupported binary knowledge file: {path.name}")
        text = _decode_text(raw)
        return [ParsedSection(path.stem, "document", text)]

    def _parse_json(self, path: Path) -> List[ParsedSection]:
        data = json.loads(_decode_text(path.read_bytes()))
        text = json.dumps(data, ensure_ascii=False, indent=2)
        return [ParsedSection(path.stem, "document", text)]

    def _parse_jsonl(self, path: Path) -> List[ParsedSection]:
        sections = []
        for index, line in enumerate(
            _decode_text(path.read_bytes()).splitlines(),
            start=1,
        ):
            if not line.strip():
                continue
            data = json.loads(line)
            sections.append(
                ParsedSection(
                    f"记录 {index}",
                    f"line {index}",
                    json.dumps(data, ensure_ascii=False, indent=2),
                )
            )
        return sections

    def _parse_delimited(
        self,
        path: Path,
        suffix: str,
    ) -> List[ParsedSection]:
        text = _decode_text(path.read_bytes())
        dialect = "\t" if suffix == ".tsv" else ","
        rows = csv.reader(io.StringIO(text), delimiter=dialect)
        rendered = ["\t".join(cell.strip() for cell in row) for row in rows]
        return [ParsedSection(path.stem, "table", "\n".join(rendered))]

    def _parse_xml(self, path: Path) -> List[ParsedSection]:
        root = ElementTree.fromstring(path.read_bytes())
        text = "\n".join(part.strip() for part in root.itertext() if part.strip())
        return [ParsedSection(root.tag, "document", text)]

    def _html_text(self, path: Path) -> str:
        parser = _VisibleHTMLParser()
        parser.feed(_decode_text(path.read_bytes()))
        return parser.text()

    def _parse_pdf(self, path: Path) -> List[ParsedSection]:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise ValueError(
                "PDF parsing requires the 'pypdf' package"
            ) from exc
        reader = PdfReader(str(path))
        return [
            ParsedSection(
                f"第 {index} 页",
                f"page {index}",
                page.extract_text() or "",
            )
            for index, page in enumerate(reader.pages, start=1)
        ]

    def _parse_docx(self, path: Path) -> List[ParsedSection]:
        try:
            from docx import Document
        except ImportError as exc:
            raise ValueError(
                "DOCX parsing requires the 'python-docx' package"
            ) from exc
        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs]
        for table_index, table in enumerate(document.tables, start=1):
            parts.append(f"\n[表格 {table_index}]")
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return [ParsedSection(path.stem, "document", "\n".join(parts))]

    def _parse_pptx(self, path: Path) -> List[ParsedSection]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValueError(
                "PPTX parsing requires the 'python-pptx' package"
            ) from exc
        presentation = Presentation(str(path))
        sections = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    parts.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.append("\t".join(cell.text for cell in row.cells))
            sections.append(
                ParsedSection(
                    f"幻灯片 {index}",
                    f"slide {index}",
                    "\n".join(parts),
                )
            )
        return sections

    def _parse_xlsx(self, path: Path) -> List[ParsedSection]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ValueError(
                "XLSX parsing requires the 'openpyxl' package"
            ) from exc
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        sections = []
        try:
            for sheet in workbook.worksheets:
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    values = [
                        "" if value is None else str(value)
                        for value in row
                    ]
                    if any(values):
                        rows.append("\t".join(values))
                sections.append(
                    ParsedSection(sheet.title, f"sheet {sheet.title}", "\n".join(rows))
                )
        finally:
            workbook.close()
        return sections

    def _parse_odt(self, path: Path) -> List[ParsedSection]:
        with zipfile.ZipFile(path) as archive:
            raw = archive.read("content.xml")
        root = ElementTree.fromstring(raw)
        text = "\n".join(part.strip() for part in root.itertext() if part.strip())
        return [ParsedSection(path.stem, "document", text)]

    def _parse_epub(self, path: Path) -> List[ParsedSection]:
        sections = []
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                name
                for name in archive.namelist()
                if Path(name).suffix.lower() in {".html", ".htm", ".xhtml"}
            )
            for name in names:
                parser = _VisibleHTMLParser()
                parser.feed(_decode_text(archive.read(name)))
                sections.append(ParsedSection(Path(name).stem, name, parser.text()))
        return sections

    @staticmethod
    def _bounded_sections(
        sections: Iterable[ParsedSection],
    ) -> List[ParsedSection]:
        remaining = MAX_EXTRACTED_CHARS
        bounded = []
        for section in sections:
            text = _normalize_text(section.text)
            if not text or remaining <= 0:
                continue
            text = text[:remaining]
            remaining -= len(text)
            bounded.append(
                ParsedSection(
                    heading=_normalize_text(section.heading)[:500],
                    location=_normalize_text(section.location)[:500],
                    text=text,
                )
            )
        return bounded


class KnowledgeBase:
    """SQLite-backed local project knowledge base with FTS5 retrieval."""

    def __init__(
        self,
        workspace: Workspace,
        parser: Optional[DocumentParser] = None,
        chunk_chars: int = DEFAULT_CHUNK_CHARS,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
        vector_store: Optional[ChromaVectorStore] = None,
    ) -> None:
        if chunk_chars < 200:
            raise ValueError("chunk_chars must be at least 200")
        if chunk_overlap < 0 or chunk_overlap >= chunk_chars:
            raise ValueError("chunk_overlap must be between 0 and chunk_chars")
        self.workspace = workspace
        self.storage = ProjectStorage(workspace)
        self.root = self.storage.root / "knowledge"
        self.database_path = self.root / "knowledge.db"
        self.parser = parser or DocumentParser()
        self.chunk_chars = chunk_chars
        self.chunk_overlap = chunk_overlap
        self.vector_store = vector_store or ChromaVectorStore.from_env(workspace)
        self.vector_error = ""

    def ingest(
        self,
        source: Path,
        *,
        source_name: Optional[str] = None,
        source_identity: Optional[str] = None,
    ) -> KnowledgeDocument:
        source = source.expanduser().resolve()
        sections = self.parser.parse(source)
        raw_hash = _hash_file(source)
        display_name = source_name or source.name
        identity = source_identity or str(source)
        document_id = hashlib.sha256(identity.encode("utf-8")).hexdigest()[:16]
        chunks = self._make_chunks(sections)
        imported_at = datetime.now(timezone.utc).isoformat()
        source_type = Path(display_name).suffix.lower() or "[no extension]"

        with self._connect() as connection:
            existing = connection.execute(
                "SELECT id FROM documents WHERE id = ?",
                (document_id,),
            ).fetchone()
            if existing:
                connection.execute(
                    "DELETE FROM chunks_fts WHERE document_id = ?",
                    (document_id,),
                )
                connection.execute(
                    "DELETE FROM chunks WHERE document_id = ?",
                    (document_id,),
                )
            connection.execute(
                """
                INSERT INTO documents (
                    id, source_name, source_path, source_type, content_hash,
                    chunk_count, imported_at
                ) VALUES (?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(id) DO UPDATE SET
                    source_name = excluded.source_name,
                    source_path = excluded.source_path,
                    source_type = excluded.source_type,
                    content_hash = excluded.content_hash,
                    chunk_count = excluded.chunk_count,
                    imported_at = excluded.imported_at
                """,
                (
                    document_id,
                    display_name,
                    str(source),
                    source_type,
                    raw_hash,
                    len(chunks),
                    imported_at,
                ),
            )
            for chunk_index, section in enumerate(chunks, start=1):
                cursor = connection.execute(
                    """
                    INSERT INTO chunks (
                        document_id, chunk_index, heading, location, content
                    ) VALUES (?, ?, ?, ?, ?)
                    """,
                    (
                        document_id,
                        chunk_index,
                        section.heading,
                        section.location,
                        section.text,
                    ),
                )
                chunk_id = cursor.lastrowid
                connection.execute(
                    """
                    INSERT INTO chunks_fts (
                        document_id, chunk_id, terms, content
                    ) VALUES (?, ?, ?, ?)
                    """,
                    (
                        document_id,
                        chunk_id,
                        " ".join(
                            sorted(
                                _terms(
                                    " ".join(
                                        [
                                            display_name,
                                            section.heading,
                                            section.location,
                                            section.text,
                                        ]
                                    )
                                )
                            )
                        ),
                        section.text,
                    ),
                )

        try:
            self.vector_store.replace_group(
                "knowledge_chunks",
                document_id,
                [
                    VectorRecord(
                        id=content_fingerprint(
                            "knowledge",
                            document_id,
                            str(chunk_index),
                        ),
                        text=section.text,
                        content_hash=content_fingerprint(
                            raw_hash,
                            str(chunk_index),
                            section.text,
                        ),
                        metadata={
                            "group": document_id,
                            "document_id": document_id,
                            "source_name": display_name,
                            "chunk_index": chunk_index,
                            "heading": section.heading,
                            "location": section.location,
                        },
                    )
                    for chunk_index, section in enumerate(chunks, start=1)
                ],
            )
            self.vector_error = ""
        except Exception as exc:
            self.vector_error = f"{type(exc).__name__}: {exc}"[:2_000]

        return KnowledgeDocument(
            document_id=document_id,
            source_name=display_name,
            source_type=source_type,
            content_hash=raw_hash,
            chunk_count=len(chunks),
            imported_at=imported_at,
        )

    def ingest_many(self, sources: Sequence[Path]) -> List[KnowledgeDocument]:
        return [self.ingest(source) for source in sources]

    def list_documents(self) -> List[KnowledgeDocument]:
        if not self.database_path.exists():
            return []
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT id, source_name, source_type, content_hash,
                       chunk_count, imported_at
                FROM documents
                ORDER BY imported_at DESC, source_name
                """
            ).fetchall()
        return [
            KnowledgeDocument(
                document_id=row["id"],
                source_name=row["source_name"],
                source_type=row["source_type"],
                content_hash=row["content_hash"],
                chunk_count=row["chunk_count"],
                imported_at=row["imported_at"],
            )
            for row in rows
        ]

    def _search_keyword(
        self,
        query: str,
        limit: int = 5,
    ) -> List[KnowledgeHit]:
        if not isinstance(query, str) or not query.strip():
            raise ValueError("knowledge query must be a non-empty string")
        if not isinstance(limit, int) or not 1 <= limit <= 20:
            raise ValueError("knowledge search limit must be from 1 to 20")
        if not self.database_path.exists():
            return []
        terms = sorted(_terms(query), key=lambda term: (-len(term), term))[:128]
        if not terms:
            return []
        expression = " OR ".join(
            f'"{term.replace(chr(34), chr(34) * 2)}"' for term in terms
        )
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT c.document_id, d.source_name, c.chunk_index,
                       c.heading, c.location, c.content,
                       bm25(chunks_fts, 0.0, 0.0, 4.0, 1.0) AS rank
                FROM chunks_fts
                JOIN chunks c ON c.id = chunks_fts.chunk_id
                JOIN documents d ON d.id = c.document_id
                WHERE chunks_fts MATCH ?
                ORDER BY rank, c.document_id, c.chunk_index
                LIMIT ?
                """,
                (expression, limit),
            ).fetchall()
        return [
            KnowledgeHit(
                document_id=row["document_id"],
                source_name=row["source_name"],
                chunk_index=row["chunk_index"],
                heading=row["heading"],
                location=row["location"],
                content=row["content"],
                score=round(-float(row["rank"]), 6),
            )
            for row in rows
        ]

    def search(self, query: str, limit: int = 5) -> List[KnowledgeHit]:
        """Fuse keyword BM25 and Chroma vector matches."""
        keyword_hits = self._search_keyword(
            query,
            min(20, max(limit * 3, limit)),
        )
        try:
            vector_hits = self.vector_store.query(
                "knowledge_chunks",
                query,
                min(20, max(limit * 3, limit)),
            )
            self.vector_error = ""
        except Exception as exc:
            self.vector_error = f"{type(exc).__name__}: {exc}"[:2_000]
            vector_hits = []
        keyword_ids = [
            f"{hit.document_id}:{hit.chunk_index}" for hit in keyword_hits
        ]
        vector_ids = [
            f"{hit.metadata.get('document_id', '')}:"
            f"{int(hit.metadata.get('chunk_index', 0))}"
            for hit in vector_hits
        ]
        fused = reciprocal_rank_fusion([keyword_ids, vector_ids])
        candidates = {
            f"{hit.document_id}:{hit.chunk_index}": hit
            for hit in keyword_hits
        }
        for hit in vector_hits:
            document_id = str(hit.metadata.get("document_id", ""))
            chunk_index = int(hit.metadata.get("chunk_index", 0))
            key = f"{document_id}:{chunk_index}"
            if key in candidates or not document_id or chunk_index < 1:
                continue
            candidates[key] = KnowledgeHit(
                document_id=document_id,
                source_name=str(hit.metadata.get("source_name", "")),
                chunk_index=chunk_index,
                heading=str(hit.metadata.get("heading", "")),
                location=str(hit.metadata.get("location", "")),
                content=hit.text,
                score=0.0,
            )
        ordered = sorted(
            candidates.items(),
            key=lambda item: (-fused.get(item[0], 0.0), item[0]),
        )
        return [
            KnowledgeHit(
                document_id=hit.document_id,
                source_name=hit.source_name,
                chunk_index=hit.chunk_index,
                heading=hit.heading,
                location=hit.location,
                content=hit.content,
                score=round(fused.get(key, 0.0), 8),
            )
            for key, hit in ordered[:limit]
        ]

    def read_chunk(self, document_id: str, chunk_index: int) -> KnowledgeHit:
        self._validate_document_id(document_id)
        if not isinstance(chunk_index, int) or chunk_index < 1:
            raise ValueError("chunk_index must be a positive integer")
        if not self.database_path.exists():
            raise ValueError(f"knowledge document does not exist: {document_id}")
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT c.document_id, d.source_name, c.chunk_index,
                       c.heading, c.location, c.content
                FROM chunks c
                JOIN documents d ON d.id = c.document_id
                WHERE c.document_id = ? AND c.chunk_index = ?
                """,
                (document_id, chunk_index),
            ).fetchone()
        if row is None:
            raise ValueError(
                f"knowledge chunk does not exist: {document_id}#{chunk_index}"
            )
        return KnowledgeHit(
            document_id=row["document_id"],
            source_name=row["source_name"],
            chunk_index=row["chunk_index"],
            heading=row["heading"],
            location=row["location"],
            content=row["content"],
            score=0.0,
        )

    def remove(self, document_id: str) -> bool:
        self._validate_document_id(document_id)
        if not self.database_path.exists():
            return False
        with self._connect() as connection:
            exists = connection.execute(
                "SELECT 1 FROM documents WHERE id = ?",
                (document_id,),
            ).fetchone()
            if not exists:
                return False
            connection.execute(
                "DELETE FROM chunks_fts WHERE document_id = ?",
                (document_id,),
            )
            connection.execute(
                "DELETE FROM documents WHERE id = ?",
                (document_id,),
            )
        try:
            self.vector_store.remove_group("knowledge_chunks", document_id)
            self.vector_error = ""
        except Exception as exc:
            self.vector_error = f"{type(exc).__name__}: {exc}"[:2_000]
        return True

    def _make_chunks(
        self,
        sections: Sequence[ParsedSection],
    ) -> List[ParsedSection]:
        chunks = []
        for section in sections:
            for text in _split_text(
                section.text,
                self.chunk_chars,
                self.chunk_overlap,
            ):
                chunks.append(
                    ParsedSection(section.heading, section.location, text)
                )
        return chunks

    def _connect(self) -> sqlite3.Connection:
        self._ensure_storage_path()
        self.root.mkdir(parents=True, exist_ok=True)
        self._ensure_storage_path()
        connection = sqlite3.connect(str(self.database_path))
        connection.row_factory = sqlite3.Row
        connection.execute("PRAGMA foreign_keys = ON")
        connection.executescript(
            """
            CREATE TABLE IF NOT EXISTS metadata (
                key TEXT PRIMARY KEY,
                value TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                source_name TEXT NOT NULL,
                source_path TEXT NOT NULL,
                source_type TEXT NOT NULL,
                content_hash TEXT NOT NULL,
                chunk_count INTEGER NOT NULL,
                imported_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS chunks (
                id INTEGER PRIMARY KEY,
                document_id TEXT NOT NULL REFERENCES documents(id)
                    ON DELETE CASCADE,
                chunk_index INTEGER NOT NULL,
                heading TEXT NOT NULL,
                location TEXT NOT NULL,
                content TEXT NOT NULL,
                UNIQUE(document_id, chunk_index)
            );
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks_fts USING fts5(
                document_id UNINDEXED,
                chunk_id UNINDEXED,
                terms,
                content,
                tokenize = 'unicode61 remove_diacritics 2'
            );
            """
        )
        connection.execute(
            """
            INSERT INTO metadata(key, value) VALUES('version', ?)
            ON CONFLICT(key) DO UPDATE SET value = excluded.value
            """,
            (str(KNOWLEDGE_VERSION),),
        )
        return connection

    def _ensure_storage_path(self) -> None:
        self.storage.ensure_path(self.database_path, "knowledge")

    @staticmethod
    def _validate_document_id(document_id: str) -> None:
        if not isinstance(document_id, str) or not VALID_DOCUMENT_ID.fullmatch(
            document_id
        ):
            raise ValueError(f"invalid knowledge document id: {document_id}")


def document_to_dict(document: KnowledgeDocument) -> Dict[str, Any]:
    """Serialize public document metadata without leaking its source path."""
    return asdict(document)


def hit_to_dict(hit: KnowledgeHit, include_content: bool = True) -> Dict[str, Any]:
    """Serialize a search hit with a stable citation."""
    data = {
        "citation": hit.citation,
        "document_id": hit.document_id,
        "source_name": hit.source_name,
        "chunk_index": hit.chunk_index,
        "heading": hit.heading,
        "location": hit.location,
        "score": hit.score,
    }
    if include_content:
        data["content"] = hit.content
    return data


def is_supported_document(path: Path) -> bool:
    """Return whether a directory import should attempt to parse this file."""
    suffix = path.suffix.lower()
    return (
        suffix in TEXT_EXTENSIONS
        or suffix in STRUCTURED_EXTENSIONS
        or path.name.lower()
        in {"dockerfile", "gemfile", "makefile", "readme", "license"}
    )


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("unable to decode text document")


def _normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _split_text(text: str, chunk_chars: int, overlap: int) -> List[str]:
    text = _normalize_text(text)
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_chars, len(text))
        if end < len(text):
            minimum = start + int(chunk_chars * 0.6)
            boundaries = [
                text.rfind("\n\n", minimum, end),
                text.rfind("\n", minimum, end),
                text.rfind(" ", minimum, end),
            ]
            boundary = max(boundaries)
            if boundary > start:
                end = boundary
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        next_start = max(start + 1, end - overlap)
        while next_start < end and text[next_start].isspace():
            next_start += 1
        start = next_start
    return chunks


def _terms(text: str) -> Set[str]:
    terms: Set[str] = set()
    for chunk in re.findall(
        r"[a-zA-Z0-9_./-]+|[\u3400-\u9fff]+",
        text.lower(),
    ):
        if re.fullmatch(r"[\u3400-\u9fff]+", chunk):
            if len(chunk) <= 12:
                terms.add(chunk)
            if len(chunk) > 1:
                terms.update(
                    chunk[index : index + 2]
                    for index in range(len(chunk) - 1)
                )
        elif len(chunk) > 1:
            terms.add(chunk)
    return terms


def _hash_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file:
        while True:
            block = file.read(1024 * 1024)
            if not block:
                break
            digest.update(block)
    return digest.hexdigest()


def _validate_archive(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise ValueError(
                    f"document archive contains too many entries: {path.name}"
                )
            expanded_size = sum(entry.file_size for entry in entries)
            if expanded_size > MAX_ARCHIVE_BYTES:
                raise ValueError(
                    f"expanded document exceeds {MAX_ARCHIVE_BYTES} bytes: "
                    f"{path.name}"
                )
    except zipfile.BadZipFile as exc:
        raise ValueError(f"invalid document archive: {path.name}") from exc
