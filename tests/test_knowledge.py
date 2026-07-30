from tests import _TEST_STORAGE_HOME  # noqa: F401

import json
import os
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from types import SimpleNamespace
from unittest.mock import patch

from simple_agent.cli import _handle_knowledge_actions, build_agent
from simple_agent.knowledge import DocumentParser, KnowledgeBase
from simple_agent.memory import ContextBuilder, ProjectMemoryStore
from simple_agent.session import SessionManager
from simple_agent.tools import (
    ListKnowledgeTool,
    ReadKnowledgeTool,
    ReadOnlyCommandTool,
    RunCommandTool,
    SearchKnowledgeTool,
)
from simple_agent.workspace import Workspace


class KnowledgeBaseTests(unittest.TestCase):
    def test_ingest_search_read_list_and_remove(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "开发规范.md"
            source.write_text(
                "# API 规范\n\n所有写接口必须支持幂等键，并返回统一错误码。\n",
                encoding="utf-8",
            )
            knowledge = KnowledgeBase(
                Workspace(root),
                chunk_chars=200,
                chunk_overlap=20,
            )

            document = knowledge.ingest(source)
            hits = knowledge.search("写接口的幂等要求", limit=3)
            loaded = knowledge.read_chunk(
                hits[0].document_id,
                hits[0].chunk_index,
            )

            self.assertEqual(document.chunk_count, 1)
            self.assertEqual(len(knowledge.list_documents()), 1)
            self.assertIn("幂等键", hits[0].content)
            self.assertIn("统一错误码", loaded.content)
            self.assertTrue(hits[0].citation.startswith("knowledge:"))
            self.assertTrue(knowledge.remove(document.document_id))
            self.assertEqual(knowledge.list_documents(), [])

    def test_reimport_replaces_old_chunks_instead_of_duplicating(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "rules.txt"
            source.write_text("Use REST for public APIs.", encoding="utf-8")
            knowledge = KnowledgeBase(Workspace(root))
            first = knowledge.ingest(source)
            source.write_text("Use GraphQL for public APIs.", encoding="utf-8")

            second = knowledge.ingest(source)

            self.assertEqual(first.document_id, second.document_id)
            self.assertEqual(len(knowledge.list_documents()), 1)
            self.assertEqual(knowledge.search("REST"), [])
            self.assertIn("GraphQL", knowledge.search("GraphQL")[0].content)

    def test_large_document_is_split_into_overlapping_chunks(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "long.txt"
            source.write_text(
                "\n\n".join(
                    f"Section {index} requires validation and review."
                    for index in range(100)
                ),
                encoding="utf-8",
            )
            knowledge = KnowledgeBase(
                Workspace(root),
                chunk_chars=300,
                chunk_overlap=50,
            )

            document = knowledge.ingest(source)

            self.assertGreater(document.chunk_count, 5)
            self.assertIn("Section 75", knowledge.search("Section 75")[0].content)

    def test_storage_rejects_symlinked_knowledge_directory(self):
        with TemporaryDirectory() as directory:
            root = Path(directory) / "workspace"
            outside = Path(directory) / "outside"
            root.mkdir()
            outside.mkdir()
            internal = root / ".simple-agent"
            internal.mkdir()
            (internal / "knowledge").symlink_to(
                outside,
                target_is_directory=True,
            )
            source = root / "rules.txt"
            source.write_text("safe content", encoding="utf-8")

            with self.assertRaisesRegex(ValueError, "symbolic links"):
                KnowledgeBase(Workspace(root)).ingest(source)


class DocumentParserTests(unittest.TestCase):
    def test_parses_json_csv_html_and_unknown_utf8_text(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "rules.json").write_text(
                '{"retry": 3, "policy": "backoff"}',
                encoding="utf-8",
            )
            (root / "owners.csv").write_text(
                "module,owner\napi,platform\n",
                encoding="utf-8",
            )
            (root / "guide.html").write_text(
                "<h1>Guide</h1><script>ignore()</script><p>Use TLS</p>",
                encoding="utf-8",
            )
            (root / "NOTICE").write_text(
                "Review all migrations.",
                encoding="utf-8",
            )
            parser = DocumentParser()

            json_text = parser.parse(root / "rules.json")[0].text
            csv_text = parser.parse(root / "owners.csv")[0].text
            html_text = parser.parse(root / "guide.html")[0].text
            notice_text = parser.parse(root / "NOTICE")[0].text

            self.assertIn('"retry": 3', json_text)
            self.assertIn("api\tplatform", csv_text)
            self.assertIn("Use TLS", html_text)
            self.assertNotIn("ignore()", html_text)
            self.assertIn("migrations", notice_text)

    def test_rejects_unsupported_binary_and_legacy_office(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "data.bin").write_bytes(b"abc\x00def")
            (root / "old.doc").write_bytes(b"legacy")
            parser = DocumentParser()

            with self.assertRaisesRegex(ValueError, "binary"):
                parser.parse(root / "data.bin")
            with self.assertRaisesRegex(ValueError, "legacy Office"):
                parser.parse(root / "old.doc")

    def test_rejects_corrupt_office_archive(self):
        with TemporaryDirectory() as directory:
            path = Path(directory) / "broken.docx"
            path.write_bytes(b"not a zip archive")

            with self.assertRaisesRegex(ValueError, "invalid document archive"):
                DocumentParser().parse(path)

    def test_parses_docx_pptx_and_xlsx(self):
        try:
            from docx import Document
            from openpyxl import Workbook
            from pptx import Presentation
        except ImportError:
            self.skipTest("office document parsing dependencies unavailable")

        with TemporaryDirectory() as directory:
            root = Path(directory)
            docx_path = root / "rules.docx"
            document = Document()
            document.add_paragraph("All changes require review.")
            document.save(docx_path)

            pptx_path = root / "architecture.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "Event-driven architecture"
            presentation.save(pptx_path)

            xlsx_path = root / "owners.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Owners"
            sheet.append(["module", "owner"])
            sheet.append(["api", "platform"])
            workbook.save(xlsx_path)

            parser = DocumentParser()

            self.assertIn("require review", parser.parse(docx_path)[0].text)
            self.assertIn(
                "Event-driven architecture",
                parser.parse(pptx_path)[0].text,
            )
            self.assertIn("api\tplatform", parser.parse(xlsx_path)[0].text)


class KnowledgeContextTests(unittest.TestCase):
    def test_relevant_knowledge_is_injected_without_full_document(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "security.md"
            source.write_text(
                "认证接口必须限制失败重试次数。" + "无关附录。" * 500,
                encoding="utf-8",
            )
            workspace = Workspace(root)
            knowledge = KnowledgeBase(
                workspace,
                chunk_chars=300,
                chunk_overlap=30,
            )
            knowledge.ingest(source)
            builder = ContextBuilder(
                ProjectMemoryStore(workspace),
                knowledge_base=knowledge,
                max_knowledge_chars=2_000,
            )

            built = builder.build("实现认证接口失败重试限制")
            serialized = json.dumps(built.messages, ensure_ascii=False)

            self.assertIn("认证接口必须限制失败重试次数", serialized)
            self.assertLess(len(serialized), 5_000)
            self.assertTrue(built.knowledge_citations)

    def test_session_records_knowledge_provenance(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "style.md"
            source.write_text("数据库字段使用 snake_case。", encoding="utf-8")
            workspace = Workspace(root)
            knowledge = KnowledgeBase(workspace)
            knowledge.ingest(source)

            session = SessionManager(
                ProjectMemoryStore(workspace),
                knowledge_base=knowledge,
            ).start_task("新增数据库字段")

            self.assertTrue(session.knowledge_citations)
            self.assertIn("knowledge:", session.knowledge_citations[0])

    def test_irrelevant_knowledge_is_not_injected(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "database.md"
            source.write_text(
                "数据库字段必须使用 snake_case。",
                encoding="utf-8",
            )
            workspace = Workspace(root)
            knowledge = KnowledgeBase(workspace)
            knowledge.ingest(source)

            built = ContextBuilder(
                ProjectMemoryStore(workspace),
                knowledge_base=knowledge,
            ).build("调整前端按钮颜色")

            self.assertEqual(len(built.messages), 1)
            self.assertIn("调整前端按钮颜色", built.messages[0]["content"])
            self.assertEqual(built.knowledge_citations, [])


class KnowledgeToolTests(unittest.TestCase):
    def test_tools_expose_bounded_knowledge_retrieval(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "testing.md"
            source.write_text("支付模块必须运行集成测试。", encoding="utf-8")
            knowledge = KnowledgeBase(Workspace(root))
            document = knowledge.ingest(source)

            listing = ListKnowledgeTool(knowledge).execute({})
            searched = SearchKnowledgeTool(knowledge).execute(
                {"query": "支付测试"}
            )
            loaded = ReadKnowledgeTool(knowledge).execute(
                {
                    "document_id": document.document_id,
                    "chunk_index": 1,
                }
            )

            self.assertIn("testing.md", listing)
            self.assertIn("支付模块", searched)
            self.assertIn("knowledge:", searched)
            self.assertIn("集成测试", loaded)


class KnowledgeCliTests(unittest.TestCase):
    def test_cli_actions_import_list_and_remove_without_llm(self):
        with TemporaryDirectory() as directory:
            root = Path(directory)
            docs = root / "docs"
            docs.mkdir()
            source = docs / "rules.md"
            source.write_text("提交前必须运行测试。", encoding="utf-8")
            knowledge = KnowledgeBase(Workspace(root))
            import_args = SimpleNamespace(
                knowledge_file=[],
                knowledge_dir=[docs],
                remove_knowledge=[],
                list_knowledge=True,
            )

            output = _handle_knowledge_actions(import_args, knowledge)
            document = knowledge.list_documents()[0]
            remove_args = SimpleNamespace(
                knowledge_file=[],
                knowledge_dir=[],
                remove_knowledge=[document.document_id],
                list_knowledge=False,
            )
            removed = _handle_knowledge_actions(remove_args, knowledge)

            self.assertIn("已导入知识文档", output)
            self.assertIn("rules.md", output)
            self.assertIn("已删除知识文档", removed)
            self.assertEqual(knowledge.list_documents(), [])

    def test_agent_registers_all_knowledge_retrieval_tools(self):
        with TemporaryDirectory() as directory:
            environment = {
                "LLM_MODEL": "test-model",
                "LLM_BASE_URL": "https://example.com/v1",
                "LLM_API_KEY": "test-key",
            }
            with patch.dict(os.environ, environment):
                agent = build_agent(Path(directory))

            names = {
                definition["function"]["name"]
                for definition in agent.tools.definitions
            }
            self.assertIn("search_knowledge", names)
            self.assertIn("read_knowledge", names)
            self.assertIn("list_knowledge", names)
            self.assertIn("project_overview", names)
            self.assertIn("query_project_index", names)
            self.assertIn("search_symbols", names)
            self.assertIn("find_references", names)
            self.assertIn("dependency_graph", names)
            self.assertIn("refresh_project_index", names)
            self.assertIn("project_graph_overview", names)
            self.assertIn("query_file_profiles", names)
            self.assertIn("file_profile", names)
            self.assertIn("query_project_graph", names)
            self.assertIn("impact_analysis", names)
            self.assertIn("refresh_project_graph", names)
            self.assertIsInstance(
                agent.tools._tools["run_command"],
                RunCommandTool,
            )
            self.assertNotIsInstance(
                agent.tools._tools["run_command"],
                ReadOnlyCommandTool,
            )
            self.assertIsInstance(
                agent.reviewer.agent.tools._tools["run_command"],
                ReadOnlyCommandTool,
            )


if __name__ == "__main__":
    unittest.main()
